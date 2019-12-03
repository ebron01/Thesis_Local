from pptx import Presentation
from pptx.util import Inches
import json
import csv
import sys
reload(sys)
sys.setdefaultencoding('utf8')
import pdb
#this loads resulting json of rmaccompare.py script. For details refer to it.
with open('query.json', 'r') as f:
    best_diff = json.load(f)

#these are splits in ActivityNet detail json to get captions for videos.
split = {'train': 'train', 'val_1': 'val_1', 'val_2': 'val_2'}

#this part unites infos of splits to a dict
vdata = {}
vkeys = []
for key in split.keys():
    filename = '/data/shared/ConceptualCaptions/keras_rmac/data/activitynet_annotations/ActivityNet_%s.json'%key
    with open(filename, 'r') as f:
        data = json.load(f)
        vdata.update({key: data})
#this part unites ids of splits to a list
for key in split.keys():
    for k in vdata[key].keys():
        vkeys.append(k)

#information about ActivityNet dataset
actnet_filename = '/data/shared/ConceptualCaptions/keras_rmac/data/activitynet_annotations/actnet_video_details.json'
with open(actnet_filename, 'r') as f:
    actnet = json.load(f)

#information about Conceptual Caption dataset
#id naming convention for ids starts from 0.jpg so index will be sufficent to get images.
file = '/data/shared/ConceptualCaptions/keras_rmac/../downloads_26Oct/Train_GCC-training.tsv'
f = open(file, 'r')
print ('file read is successful!')
d = csv.reader(f, delimiter="\t")
data = []
for i in d:
    data.append(i)

#Conceptual Caption datasets directory folder in servers.
server_concap_image_pth = '/data/shared/ConceptualCaptions/downloads_26Oct/'
server_activity_image_path = '/data/shared/ConceptualCaptions/keras_rmac/ActVideos/frames/'

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

#keys of activitynet videos are ordered
for act_key in sorted(best_diff.keys()):
    #keys of conceptual datasets closes results are reordered according to their cosine similarity asc.
    for concat_key in sorted(best_diff[act_key].items(), key=lambda kv: kv[1]):

        slide = prs.slides.add_slide(blank_slide_layout)
        video_id = str(act_key.rsplit('_', 2)[0])
        event_count = int(act_key.rsplit('_', 2)[1]) -1 # extract -1 from event count because while middle frames of events are named, event counts starts from 1. Event index for first event in a video starts from 1.
        frame_id = str(act_key.rsplit('_', 2)[2])
        split = actnet[video_id]['split']
        sentence = vdata[split][video_id]['sentences'][event_count]

        left = top = width = height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "ActivityNet"

        # this part adds activity pictures
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(2.5)
        img_path_act = server_activity_image_path + video_id + '/' + frame_id + '.jpg'
        pic = slide.shapes.add_picture(img_path_act, left, top, height=height)

        #this part adds activity text
        left =Inches(0.5)
        top = Inches(5)
        width = height = Inches(4)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.text = act_key + '.jpg' + '\n' + sentence
        # p = tf.add_paragraph()
        # p.text = sentence

        # this part adds concap pictures
        top = width = height = Inches(0.5)
        left = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Conceptual Captions"

        # this part adds concap pictures
        left = Inches(5)
        top = Inches(1.5)
        height = Inches(2.5)

        if len(concat_key[0]) < 6:
            server_concap_image_path = server_concap_image_pth + 'images0_100/'
        elif str(concat_key[0][0]) == '1':
            server_concap_image_path = server_concap_image_pth + 'images100-200/'
        else:
            # pdb.set_trace()
            server_concap_image_path = server_concap_image_pth + 'images%d00_%d00/'%((int(concat_key[0][0])), (int(concat_key[0][0])+1))

        img_path = server_concap_image_path + concat_key[0] + '.jpg'
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

        # this part adds concap text
        width = height = Inches(4.5)
        left = top = Inches(5)
        concapBox = slide.shapes.add_textbox(left, top, width, height)
        cf = concapBox.text_frame
        cf.word_wrap = True
        concat_sentence = data[int(concat_key[0])][0]
        cf.text = concat_key[0] + '.jpg  ' + 'Distance:' + str(best_diff[act_key][concat_key[0]]) + '\n' + str(concat_sentence)
        # cf.text = concat_key + '.jpg  ' + 'Cosine Similarty:' + str(best_10_diff[act_key][concat_key])

        # p = tf.add_paragraph()
        # p.text = concat_key + '.jpg  ' + 'Cosine Similarty:' + str(best_10_diff[act_key][concat_key])
        # p.font.bold = True

prs.save('server_imgsim_best_30_0_1.pptx')
